# scripts/build_deck.py
# Generates the Core-Roborate demo + roadmap deck as a .pptx.
# Brand-aligned with CLAUDE.md tokens: dark slate, indigo accent, blue/amber semantics.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── brand tokens (from CLAUDE.md) ─────────────────────────────
SIDEBAR_BG   = RGBColor(0x0F, 0x11, 0x17)
CONTENT_BG   = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BORDER  = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
TEXT_SECOND  = RGBColor(0x47, 0x55, 0x69)
TEXT_MUTED   = RGBColor(0x94, 0xA3, 0xB8)
ACCENT       = RGBColor(0x63, 0x66, 0xF1)  # indigo-500
ACCENT_DK    = RGBColor(0x4F, 0x46, 0xE5)
VERIFIED     = RGBColor(0x3B, 0x82, 0xF6)  # blue
AI_AMBER     = RGBColor(0xF5, 0x9E, 0x0B)  # amber
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS      = RGBColor(0x10, 0xB9, 0x81)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide(bg=CONTENT_BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)  # rectangle
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_para(p, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             italic=False, space_after=6, font=FONT):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.runs[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return p


def add_para(tf, text, size, color, **kw):
    p = tf.add_paragraph()
    return set_para(p, text, size, color, **kw)


def accent_bar(slide, color=ACCENT, top=Inches(0.0), height=Inches(0.12)):
    bar = slide.shapes.add_shape(1, 0, top, SW, height)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False
    return bar


def kicker(slide, text, color=ACCENT, top=Inches(0.55)):
    tf = textbox(slide, Inches(0.7), top, Inches(11.9), Inches(0.4))
    set_para(tf.paragraphs[0], text.upper(), 13, color, bold=True)
    tf.paragraphs[0].runs[0].font.spacing = Pt(2)


def title(slide, text, top=Inches(0.95), color=TEXT_PRIMARY, size=34):
    tf = textbox(slide, Inches(0.7), top, Inches(11.9), Inches(1.0))
    set_para(tf.paragraphs[0], text, size, color, bold=True)


def footer(slide, idx, total, label="Core-Roborate — Confidential"):
    tf = textbox(slide, Inches(0.7), Inches(7.02), Inches(11.9), Inches(0.35))
    p = tf.paragraphs[0]
    set_para(p, label, 9, TEXT_MUTED)
    tf2 = textbox(slide, Inches(11.2), Inches(7.02), Inches(1.4), Inches(0.35))
    set_para(tf2.paragraphs[0], f"{idx} / {total}", 9, TEXT_MUTED, align=PP_ALIGN.RIGHT)


def bullets(slide, items, left=Inches(0.8), top=Inches(2.1),
            width=Inches(11.7), size=17, gap=10, color=TEXT_SECOND):
    tf = textbox(slide, left, top, width, Inches(4.4))
    first = True
    for it in items:
        txt, lvl = (it if isinstance(it, tuple) else (it, 0))
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        bullet = "•   " if lvl == 0 else "–   "
        c = color if lvl == 0 else TEXT_MUTED
        sz = size if lvl == 0 else size - 2
        set_para(p, bullet + txt, sz, c, space_after=gap)
        p.level = lvl
    return tf


def card(slide, l, t, w, h, fill=WHITE, border=CARD_BORDER, line_w=1.0):
    box = slide.shapes.add_shape(5, l, t, w, h)  # rounded rect
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = border; box.line.width = Pt(line_w)
    box.shadow.inherit = False
    try:
        box.adjustments[0] = 0.06
    except Exception:
        pass
    return box


def visual_note(slide, text):
    tf = textbox(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.6))
    p = tf.paragraphs[0]
    set_para(p, "VISUAL:  " + text, 11, ACCENT, italic=True)


TOTAL = 18

# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
s = add_slide(SIDEBAR_BG)
accent_bar(s, ACCENT, top=Inches(0), height=Inches(0.18))
tf = textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.2))
set_para(tf.paragraphs[0], "Core-Roborate", 60, WHITE, bold=True)
add_para(tf, "The System of Record for Workforce Truth", 24, RGBColor(0xC7,0xD2,0xFE), space_after=4)
add_para(tf, "Current Build Demo  +  Product Roadmap", 16, TEXT_MUTED)
tf2 = textbox(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5))
set_para(tf2.paragraphs[0], "Presenter  •  Date  •  Confidential", 12, TEXT_MUTED)
# shield + sparkle motif
chip = card(s, Inches(10.4), Inches(2.6), Inches(2.0), Inches(0.6), fill=RGBColor(0x1E,0x24,0x33), border=VERIFIED)
ctf = chip.text_frame; ctf.word_wrap=True
set_para(ctf.paragraphs[0], "🛡  Verified", 14, VERIFIED, bold=True, align=PP_ALIGN.CENTER)
chip2 = card(s, Inches(10.4), Inches(3.35), Inches(2.0), Inches(0.6), fill=RGBColor(0x2A,0x22,0x10), border=AI_AMBER)
set_para(chip2.text_frame.paragraphs[0], "✦  AI Estimate", 14, AI_AMBER, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "The Problem")
title(s, "Every talent tool blurs opinion and fact.")
bullets(s, [
    "Résumés are self-reported and impossible to verify.",
    "Performance data mixes hard evidence with manager gut-feel.",
    "AI tools confidently present guesses as if they were truth.",
    "The result: nobody can fully trust workforce data today.",
], top=Inches(2.2), size=19, gap=16)
visual_note(s, "Messy 'blended data' graphic — a muddy purple blob mixing fact + opinion.")
footer(s, 2, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — THE INSIGHT
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "The Insight — Our Wedge")
title(s, "We made the distinction architectural.")
# two cards
c1 = card(s, Inches(0.8), Inches(2.2), Inches(5.6), Inches(3.4), fill=RGBColor(0xEF,0xF6,0xFF), border=VERIFIED, line_w=1.5)
t1 = c1.text_frame; t1.word_wrap=True; t1.margin_left=Inches(0.3); t1.margin_top=Inches(0.3)
set_para(t1.paragraphs[0], "🛡  Verified Facts", 22, VERIFIED, bold=True)
add_para(t1, "Attested by real humans", 16, TEXT_PRIMARY, space_after=4)
add_para(t1, "Stored in separate verified_* tables", 14, TEXT_SECOND, space_after=4)
add_para(t1, "Shield icon · blue · 5-level confidence", 14, TEXT_SECOND)
c2 = card(s, Inches(6.9), Inches(2.2), Inches(5.6), Inches(3.4), fill=RGBColor(0xFF,0xFB,0xEB), border=AI_AMBER, line_w=1.5)
t2 = c2.text_frame; t2.word_wrap=True; t2.margin_left=Inches(0.3); t2.margin_top=Inches(0.3)
set_para(t2.paragraphs[0], "✦  AI Inferences", 22, AI_AMBER, bold=True)
add_para(t2, "Model-generated estimates", 16, TEXT_PRIMARY, space_after=4)
add_para(t2, "Stored in separate ai_inference_* tables", 14, TEXT_SECOND, space_after=4)
add_para(t2, "Sparkle icon · amber · always advisory", 14, TEXT_SECOND)
tf = textbox(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6))
set_para(tf.paragraphs[0], "Enforced in the database — not just a UI choice. This is the moat.", 16, ACCENT_DK, bold=True, align=PP_ALIGN.CENTER)
footer(s, 3, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT CREDENTIA IS
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "What Core-Roborate Is")
title(s, "Enterprise workforce verification & talent intelligence.")
bullets(s, [
    "Multi-tenant SaaS — each company is fully isolated.",
    "Six roles, each with a purpose-built dashboard.",
    "Turns everyday work into provable, portable credentials.",
], top=Inches(2.2), size=19, gap=14)
# role pyramid (simple stacked bars)
roles = ["superadmin","admin","executive","manager","employee","former employee"]
top0 = Inches(4.4); h=Inches(0.4)
for i, r in enumerate(roles):
    w = Inches(3.0 + i*1.4)
    bar = card(s, Inches(0.8), Emu(int(top0)+i*int(Inches(0.45))), w, h, fill=ACCENT if i==0 else RGBColor(0xE0,0xE7,0xFF), border=ACCENT)
    set_para(bar.text_frame.paragraphs[0], r, 12, WHITE if i==0 else ACCENT_DK, bold=True, align=PP_ALIGN.LEFT)
    bar.text_frame.margin_left=Inches(0.2)
visual_note(s, "Role hierarchy pyramid, superadmin at top narrowing to former employee.")
footer(s, 4, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — SECTION DIVIDER
# ════════════════════════════════════════════════════════════
s = add_slide(SIDEBAR_BG)
accent_bar(s, ACCENT, top=Inches(3.5), height=Inches(0.06))
tf = textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.0), anchor=MSO_ANCHOR.MIDDLE)
set_para(tf.paragraphs[0], "Part 1 — What We've Built", 40, WHITE, bold=True)
add_para(tf, "The platform is live. Here's what works today.", 20, TEXT_MUTED)
footer(s, 5, TOTAL, "Core-Roborate — Current Build")

# ════════════════════════════════════════════════════════════
# SLIDE 6 — PLATFORM FOUNDATION
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "Current Build · Foundation")
title(s, "Built for enterprise from day one.")
bullets(s, [
    "Logically isolated multi-tenancy — org_id + Row Level Security.",
    "Tenant lifecycle: provisioning → active → suspended.",
    "Forward-compatible dedicated-isolation mode in schema.",
    "Tamper-evident, hash-chained audit trail on every action.",
], top=Inches(2.2), size=18, gap=14)
warn = card(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.7), fill=RGBColor(0xFF,0xFB,0xEB), border=AI_AMBER)
set_para(warn.text_frame.paragraphs[0], "Verbal note: dedicated-DB isolation is scaffolded, not yet live.", 13, RGBColor(0x92,0x6A,0x0A), italic=True)
warn.text_frame.margin_left=Inches(0.3)
footer(s, 6, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — VERIFICATION ENGINE (HERO)
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, VERIFIED)
kicker(s, "Current Build · Core Differentiator", color=VERIFIED)
title(s, "Trust, made provable.")
bullets(s, [
    "5-level verification ladder — self (L1) → multi-source (L5).",
    "Records freeze on departure — immutable, tamper-proof.",
    "Verified vs. AI separation enforced end-to-end.",
], top=Inches(2.1), size=19, gap=14)
# L1-L5 badges
labels = ["L1 Self","L2 Manager","L3 HR","L4 Company","L5 Multi-source"]
for i,l in enumerate(labels):
    intensity = 0xDB - i*0x22
    bx = Inches(0.8 + i*2.35)
    b = card(s, bx, Inches(4.6), Inches(2.15), Inches(0.85), fill=RGBColor(0x3B,0x82,0xF6) if i>=1 else RGBColor(0xBF,0xDB,0xFE), border=VERIFIED)
    set_para(b.text_frame.paragraphs[0], l, 12, WHITE if i>=1 else VERIFIED, bold=True, align=PP_ALIGN.CENTER)
visual_note(s, "Verification level badges L1–L5 + a frozen-record lock icon.")
footer(s, 7, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — EMPLOYEE & MANAGER (LIVE DEMO)
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "Current Build · LIVE DEMO")
title(s, "Daily work becomes a verified record.")
bullets(s, [
    "Employee: Achievement Vault · KPIs · Projects · Daily Pulse.",
    "Manager: swipe-to-verify deck — right approve, left deny, up clarify.",
    "Smart task delegation tied to Strategic Pillars.",
    "Promote a completed task to a Manager-Verified achievement in one tap.",
], top=Inches(2.2), size=18, gap=13)
demo = card(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.75), fill=RGBColor(0xEE,0xF2,0xFF), border=ACCENT)
set_para(demo.text_frame.paragraphs[0], "▶  LIVE DEMO: screen-record the swipe deck — this is the 'wow' moment.", 14, ACCENT_DK, bold=True)
demo.text_frame.margin_left=Inches(0.3)
footer(s, 8, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 9 — EXECUTIVE INTELLIGENCE (LIVE DEMO)
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "Current Build · LIVE DEMO")
title(s, "The whole org, at a glance.")
bullets(s, [
    "Command center: health · productivity · morale · retention · innovation.",
    "Interactive d3 org mind-map with department drill-down.",
    "Executive approval queue for high-value achievements.",
    "Everything here is built only from human-verified facts.",
], top=Inches(2.2), size=18, gap=13)
demo = card(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.75), fill=RGBColor(0xEE,0xF2,0xFF), border=ACCENT)
set_para(demo.text_frame.paragraphs[0], "▶  LIVE DEMO: open the radial org mind-map and drill into a department.", 14, ACCENT_DK, bold=True)
demo.text_frame.margin_left=Inches(0.3)
footer(s, 9, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — AI ADVISORY LAYER
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, AI_AMBER)
kicker(s, "Current Build · AI Layer", color=AI_AMBER)
title(s, "AI that earns trust — labeled, never deciding.")
bullets(s, [
    "Hybrid engine: code computes the numbers, AI writes the narrative.",
    "Compensation ranges · Value Score (0–1000) · Promotion Readiness.",
    "Retention flags — aggregate only, never names a person.",
    "Always amber, always a suggestion, always human-approved.",
], top=Inches(2.2), size=18, gap=13)
warn = card(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.75), fill=RGBColor(0xFF,0xFB,0xEB), border=AI_AMBER)
set_para(warn.text_frame.paragraphs[0], "Verbal note: runs on Claude Sonnet — wired and live, hardening for GA.", 14, RGBColor(0x92,0x6A,0x0A), italic=True)
warn.text_frame.margin_left=Inches(0.3)
footer(s, 10, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 11 — PROVISIONING / LIFECYCLE / COMPLIANCE
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "Current Build · Enterprise Plumbing")
title(s, "Provisioning, lifecycle & compliance.")
bullets(s, [
    "SSO / SCIM / Okta · CSV bulk import with dry-run validation.",
    "Portable passport — employees keep their verified record for life.",
    "Revocable public shareable profiles (verified data only).",
    "Regulator-ready PDF exports · billing & trial management.",
], top=Inches(2.2), size=18, gap=13)
warn = card(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.75), fill=RGBColor(0xFF,0xFB,0xEB), border=AI_AMBER)
set_para(warn.text_frame.paragraphs[0], "Verbal note: billing ledger is mocked — no payment processor wired yet.", 14, RGBColor(0x92,0x6A,0x0A), italic=True)
warn.text_frame.margin_left=Inches(0.3)
footer(s, 11, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 12 — SECTION DIVIDER: VISION
# ════════════════════════════════════════════════════════════
s = add_slide(SIDEBAR_BG)
accent_bar(s, ACCENT, top=Inches(3.5), height=Inches(0.06))
tf = textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.4), anchor=MSO_ANCHOR.MIDDLE)
set_para(tf.paragraphs[0], "Part 2 — Where We're Going", 40, WHITE, bold=True)
add_para(tf, "Today we verify truth inside one company.", 20, TEXT_MUTED, space_after=2)
add_para(tf, "Next, we make it portable.", 20, RGBColor(0xC7,0xD2,0xFE), bold=True)
footer(s, 12, TOTAL, "Core-Roborate — Roadmap")

# ════════════════════════════════════════════════════════════
# SLIDE 13 — PHASE 2
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "The Vision · Phase 2 — Short-Term (2–3 Quarters)")
title(s, "Deepen the moat. Monetize the data we capture.")
bullets(s, [
    ("Network: cross-company credential portability · verifier marketplace.", 0),
    ("Trust: explainable AI panels that cite the verified source · NL queries.", 0),
    ("Workflow: mobile app · succession workspace · integrations hub.", 0),
], top=Inches(2.2), size=18, gap=16)
tf = textbox(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.6))
set_para(tf.paragraphs[0], "Sequenced by leverage — quick trust wins first, network features next.", 16, ACCENT_DK, bold=True)
visual_note(s, "'More companies = more value' network-effect graphic.")
footer(s, 13, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 14 — PHASE 3
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "The Vision · Phase 3 — Long-Term (12–24 Months)")
title(s, "Become the default trust layer for the labor market.")
bullets(s, [
    "Verified Talent Graph — the 'LinkedIn that can't lie.'",
    "Verified talent marketplace · cryptographic credential anchoring.",
    "Predictive workforce planning · pay-equity audit engine.",
    "Benchmarking-as-a-service · developer / verification API.",
], top=Inches(2.2), size=18, gap=13)
visual_note(s, "Talent graph — people nodes linked across multiple company logos.")
footer(s, 14, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 15 — STRATEGIC ARC (MONEY SLIDE)
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s)
kicker(s, "The Strategic Arc")
title(s, "System → Network → Standard of Record")
cols = [
    ("Today", "Truth inside one company", "Fact vs. AI separation", RGBColor(0xE0,0xE7,0xFF)),
    ("Phase 2", "Truth portable between companies", "AI that cites its sources", RGBColor(0xC7,0xD2,0xFE)),
    ("Phase 3", "The trust layer for the market", "Prediction on a verified base", ACCENT),
]
for i,(h,a,b,fill) in enumerate(cols):
    cx = Inches(0.8 + i*4.05)
    c = card(s, cx, Inches(2.4), Inches(3.75), Inches(3.0), fill=fill, border=ACCENT)
    t = c.text_frame; t.word_wrap=True; t.margin_left=Inches(0.25); t.margin_top=Inches(0.3)
    last = i==2
    set_para(t.paragraphs[0], h, 22, WHITE if last else ACCENT_DK, bold=True)
    add_para(t, a, 15, WHITE if last else TEXT_PRIMARY, space_after=10)
    add_para(t, b, 14, RGBColor(0xE0,0xE7,0xFF) if last else TEXT_SECOND)
footer(s, 15, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSE & ASK
# ════════════════════════════════════════════════════════════
s = add_slide(SIDEBAR_BG)
accent_bar(s, ACCENT, top=Inches(0), height=Inches(0.18))
tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.5))
set_para(tf.paragraphs[0], "Core-Roborate makes workforce data\nprovable — and portable.", 34, WHITE, bold=True)
add_para(tf, "A live platform · a defensible moat · a clear path to a network.", 18, TEXT_MUTED, space_after=18)
add_para(tf, "The Ask:  [funding / pilot customers / hiring — fill per audience]", 18, RGBColor(0xC7,0xD2,0xFE), bold=True)
tf2 = textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6))
set_para(tf2.paragraphs[0], "Contact  •  Q&A", 14, TEXT_MUTED)
footer(s, 16, TOTAL, "Core-Roborate — Confidential")

# ════════════════════════════════════════════════════════════
# SLIDE 17 — APPENDIX: BUILD MATURITY
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, TEXT_MUTED)
kicker(s, "Appendix · Build Maturity", color=TEXT_SECOND)
title(s, "Demo-ready vs. GA-hardening (full transparency)")
rows = [
    ("Live", SUCCESS, "Verification engine · RLS multi-tenancy · audit trail · dashboards · provisioning · swipe-to-verify."),
    ("Wired, hardening", AI_AMBER, "AI advisory layer (Claude Sonnet) · retention flags."),
    ("Mocked", RGBColor(0xEF,0x44,0x44), "Billing / payment processor · dedicated-DB isolation mode."),
]
y = Inches(2.3)
for label, color, desc in rows:
    chip = card(s, Inches(0.8), y, Inches(2.6), Inches(0.9), fill=color, border=color)
    set_para(chip.text_frame.paragraphs[0], label, 15, WHITE, bold=True, align=PP_ALIGN.CENTER)
    tf = textbox(s, Inches(3.6), y, Inches(8.9), Inches(0.9), anchor=MSO_ANCHOR.MIDDLE)
    set_para(tf.paragraphs[0], desc, 15, TEXT_SECOND)
    y = Emu(int(y) + int(Inches(1.1)))
footer(s, 17, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 18 — APPENDIX: TECH STACK
# ════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, TEXT_MUTED)
kicker(s, "Appendix · Tech Stack", color=TEXT_SECOND)
title(s, "What it's built on")
stack = ["Next.js 16 (App Router)", "TypeScript", "Supabase (Postgres + RLS)",
         "Tailwind CSS v4", "Framer Motion", "Anthropic Claude", "Vercel"]
for i, item in enumerate(stack):
    col = i % 2
    row = i // 2
    bx = Inches(0.8 + col*6.0)
    by = Emu(int(Inches(2.3)) + row*int(Inches(0.8)))
    chip = card(s, bx, by, Inches(5.6), Inches(0.6), fill=WHITE, border=CARD_BORDER)
    set_para(chip.text_frame.paragraphs[0], item, 15, TEXT_PRIMARY, bold=True)
    chip.text_frame.margin_left=Inches(0.3)
footer(s, 18, TOTAL)

out = r"C:\Users\tyrel\credentia\Core-Roborate_Demo_Roadmap.pptx"
prs.save(out)
print("Saved:", out, "slides:", len(prs.slides._sldIdLst))
